"""
docx_to_pptx.py
===============
Génère un PowerPoint à partir de output.docx en utilisant template-ppt.potx.

Usage : python docx_to_pptx.py [input.docx] [output.pptx]

Dépendances : pip install python-pptx python-docx
"""

import os, sys, shutil, zipfile
from docx import Document
from pptx import Presentation
from pptx.util import Pt

# ── Config ───────────────────────────────────────────────────────────────────

INPUT_DOCX = "output.docx"
TEMPLATE_POTX = "template-ppt.potx"
OUTPUT_PPTX = "output.pptx"

# Layouts du template Alithya
LAYOUT_COVER = 0          # Cover-Navy 1 (titre + sous-titres)
LAYOUT_DIVIDER = 10       # Divider without number blue (titre section)
LAYOUT_CONTENT = 5        # Layout 1 column (titre + body)

# Max bullet points par slide avant split
MAX_BULLETS_PER_SLIDE = 8

# ── Helpers ──────────────────────────────────────────────────────────────────

def potx_to_pptx(potx_path):
    """Convertit .potx → .pptx temporaire (change le content type)."""
    tmp = potx_path.replace(".potx", "_tmp.pptx")
    shutil.copy2(potx_path, tmp)
    with zipfile.ZipFile(tmp, "r") as zin:
        ct = zin.read("[Content_Types].xml").decode("utf-8")
    ct = ct.replace(
        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    )
    patched = tmp + ".tmp"
    with zipfile.ZipFile(tmp, "r") as zin, zipfile.ZipFile(patched, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = ct.encode("utf-8")
            zout.writestr(item, data)
    os.replace(patched, tmp)
    return tmp


def parse_docx(path):
    """Parse le docx et retourne une liste de sections [{title, paragraphs}]."""
    doc = Document(path)
    sections = []
    # Metadata: les 2 premiers paragraphes (titre doc + date)
    meta_title, meta_date = "", ""
    for para in doc.paragraphs:
        if para.style and "Heading" in para.style.name:
            break
        t = para.text.strip()
        if not t or "table des matières" in t.lower():
            continue
        if not meta_title:
            meta_title = t
        elif not meta_date:
            meta_date = t

    current = None
    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if style == "Heading 1":
            if current:
                sections.append(current)
            current = {"title": text, "paragraphs": []}
        elif current is not None:
            current["paragraphs"].append(text)

    if current:
        sections.append(current)

    return meta_title, meta_date, sections


def add_cover(prs, title, subtitle):
    """Ajoute la slide de couverture."""
    layout = prs.slide_layouts[LAYOUT_COVER]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
        elif ph.placeholder_format.idx == 11:
            ph.text = subtitle


def add_divider(prs, title):
    """Ajoute une slide de section (divider bleu)."""
    layout = prs.slide_layouts[LAYOUT_DIVIDER]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title


def add_content_slide(prs, title, bullets):
    """Ajoute une slide de contenu avec titre + bullet points."""
    layout = prs.slide_layouts[LAYOUT_CONTENT]
    slide = prs.slides.add_slide(layout)
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
        elif ph.placeholder_format.idx == 11:
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(12)


# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    input_docx = sys.argv[1] if len(sys.argv) > 1 else INPUT_DOCX
    output_pptx = sys.argv[2] if len(sys.argv) > 2 else OUTPUT_PPTX

    if not os.path.exists(input_docx):
        print(f"❌  Fichier introuvable : {input_docx}")
        sys.exit(1)

    # 1. Convertir le template .potx → .pptx utilisable
    print(f"📄  Template : {TEMPLATE_POTX}")
    tmp_pptx = potx_to_pptx(TEMPLATE_POTX)

    # 2. Charger le template et supprimer les slides existantes
    prs = Presentation(tmp_pptx)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # 3. Parser le docx
    print(f"📖  Lecture de : {input_docx}")
    meta_title, meta_date, sections = parse_docx(input_docx)
    print(f"    Titre : {meta_title}")
    print(f"    Date  : {meta_date}")
    print(f"    Sections : {len(sections)}")

    # 4. Générer les slides
    add_cover(prs, meta_title, meta_date)

    for section in sections:
        add_divider(prs, section["title"])

        # Split les paragraphes en chunks de MAX_BULLETS_PER_SLIDE
        paras = section["paragraphs"]
        for i in range(0, len(paras), MAX_BULLETS_PER_SLIDE):
            chunk = paras[i:i + MAX_BULLETS_PER_SLIDE]
            suffix = f" ({i // MAX_BULLETS_PER_SLIDE + 1})" if len(paras) > MAX_BULLETS_PER_SLIDE else ""
            add_content_slide(prs, section["title"] + suffix, chunk)

    # 5. Sauvegarder
    prs.save(output_pptx)
    os.remove(tmp_pptx)
    print(f"✅  PowerPoint généré : {output_pptx} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
