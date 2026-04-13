"""
notion_to_pptx.py
=================
Génère un PowerPoint à partir d'une page Notion en utilisant template-ppt.potx.

Usage : python notion_to_pptx.py <notion_url> [output.pptx]

Dépendances : pip install python-pptx python-dotenv
"""

import os, re, sys, json, shutil, zipfile, urllib.request
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Pt

load_dotenv()

TEMPLATE_POTX = "template-ppt.potx"
OUTPUT_PPTX = "output.pptx"
NOTION_TOKEN = os.getenv("NOTION_TOKEN")

LAYOUT_COVER = 0
LAYOUT_DIVIDER = 10
LAYOUT_CONTENT = 5
MAX_BULLETS_PER_SLIDE = 8

# ── Notion API ───────────────────────────────────────────────────────────────

def notion_api(endpoint):
    url = f"https://api.notion.com/v1/{endpoint}"
    req = urllib.request.Request(url, headers={
        "Authorization": f"Bearer {NOTION_TOKEN}",
        "Notion-Version": "2022-06-28",
    })
    with urllib.request.urlopen(req) as resp:
        return json.loads(resp.read())


def extract_page_id(url):
    clean = url.replace("\\", "").split("?")[0].rstrip("/")
    m = re.search(r"([0-9a-f]{32})$", clean)
    if not m:
        print(f"❌  Impossible d'extraire l'ID de la page : {url}")
        sys.exit(1)
    raw = m.group(1)
    return f"{raw[:8]}-{raw[8:12]}-{raw[12:16]}-{raw[16:20]}-{raw[20:]}"


def get_all_blocks(page_id):
    blocks, cursor = [], None
    while True:
        ep = f"blocks/{page_id}/children?page_size=100"
        if cursor:
            ep += f"&start_cursor={cursor}"
        data = notion_api(ep)
        blocks.extend(data["results"])
        if not data.get("has_more"):
            break
        cursor = data["next_cursor"]
    return blocks


def get_page_title(page_id):
    page = notion_api(f"pages/{page_id}")
    for prop in page.get("properties", {}).values():
        if prop.get("type") == "title":
            return "".join(rt.get("plain_text", "") for rt in prop.get("title", []))
    return "Présentation"


def rich_text_plain(rt_list):
    return "".join(rt.get("plain_text", "") for rt in rt_list)

# ── Parse Notion blocks → sections ──────────────────────────────────────────

def parse_notion_blocks(blocks):
    """Retourne (meta_title, meta_date, sections) depuis les blocs Notion."""
    sections = []
    meta_title, meta_date = "", ""
    current = None

    for b in blocks:
        btype = b["type"]

        # Headings → nouvelle section
        if btype in ("heading_1", "heading_2", "heading_3"):
            text = rich_text_plain(b[btype]["rich_text"]).strip()
            if not text:
                continue
            if current:
                sections.append(current)
            current = {"title": text, "paragraphs": []}

        # Paragraphes, listes, callouts → contenu
        elif btype in ("paragraph", "bulleted_list_item", "numbered_list_item", "callout"):
            text = rich_text_plain(b[btype]["rich_text"]).strip()
            if not text:
                continue
            # Avant le premier heading → metadata
            if current is None:
                if not meta_title:
                    meta_title = text
                elif not meta_date:
                    meta_date = text
            else:
                current["paragraphs"].append(text)

    if current:
        sections.append(current)

    return meta_title, meta_date, sections

# ── PowerPoint generation (identique à docx_to_pptx.py) ─────────────────────

def potx_to_pptx(potx_path):
    tmp = potx_path.replace(".potx", "_tmp.pptx")
    shutil.copy2(potx_path, tmp)
    with zipfile.ZipFile(tmp, "r") as zin:
        ct = zin.read("[Content_Types].xml").decode("utf-8")
    ct = ct.replace(
        "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
    )
    patched = tmp + ".tmp"
    with zipfile.ZipFile(tmp, "r") as zin, zipfile.ZipFile(patched, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "[Content_Types].xml":
                data = ct.encode("utf-8")
            zout.writestr(item, data)
    os.replace(patched, tmp)
    return tmp


def add_cover(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_COVER])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
        elif ph.placeholder_format.idx == 11:
            ph.text = subtitle


def add_divider(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_DIVIDER])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title


def add_content_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_CONTENT])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = title
        elif ph.placeholder_format.idx == 11:
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(12)


def remove_template_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# ── Main ─────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print("Usage: python notion_to_pptx.py <notion_url> [output.pptx]")
        sys.exit(1)

    if not NOTION_TOKEN:
        print("❌  NOTION_TOKEN manquant dans .env")
        sys.exit(1)

    notion_url = sys.argv[1]
    output_pptx = sys.argv[2] if len(sys.argv) > 2 else OUTPUT_PPTX

    # 1. Récupérer les blocs Notion
    page_id = extract_page_id(notion_url)
    print(f"📄  Page Notion : {page_id}")

    page_title = get_page_title(page_id)
    blocks = get_all_blocks(page_id)
    print(f"📦  {len(blocks)} blocs récupérés")

    # 2. Parser les blocs
    meta_title, meta_date, sections = parse_notion_blocks(blocks)
    # Fallback: utiliser le titre de la page Notion si pas de meta_title
    if not meta_title:
        meta_title = page_title
    print(f"    Titre : {meta_title}")
    print(f"    Date  : {meta_date}")
    print(f"    Sections : {len(sections)}")

    # 3. Générer le PowerPoint
    tmp_pptx = potx_to_pptx(TEMPLATE_POTX)
    prs = Presentation(tmp_pptx)
    remove_template_slides(prs)

    add_cover(prs, meta_title, meta_date)

    for section in sections:
        add_divider(prs, section["title"])
        paras = section["paragraphs"]
        for i in range(0, max(1, len(paras)), MAX_BULLETS_PER_SLIDE):
            chunk = paras[i:i + MAX_BULLETS_PER_SLIDE]
            if not chunk:
                continue
            suffix = f" ({i // MAX_BULLETS_PER_SLIDE + 1})" if len(paras) > MAX_BULLETS_PER_SLIDE else ""
            add_content_slide(prs, section["title"] + suffix, chunk)

    prs.save(output_pptx)
    os.remove(tmp_pptx)
    print(f"✅  PowerPoint généré : {output_pptx} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
